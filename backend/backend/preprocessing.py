import fitz
from pdf2docx import parse
from . import Final_theme
from sumy.summarizers.lsa import LsaSummarizer
from sumy.nlp.tokenizers import Tokenizer
from sumy.parsers.plaintext import PlaintextParser
from pptx import Presentation
from docx2pdf import convert
import docx2txt
import sys
import win32com.client
import os

directory = r"C:\Users\lenovo\Desktop\FYP\backend\uploads\\"

def preprocess(iformat,oformat,nosen,fname,theme):
#IF THE INPUT FILE IS .TXT
    directory = r"C:\Users\lenovo\Desktop\FYP\backend\uploads\\"
    print("reached preprocess function")
    N=nosen
    file_type=iformat
    directory = directory + fname
    TEMP=theme

    if (file_type=="txt") :
        print("it is txt")
        file2=open(directory,'r')
        original_text=file2.read()
        file2.close()
        parser=PlaintextParser.from_string(original_text,Tokenizer('english'))
        # creating the summarizer
        lsa_summarizer=LsaSummarizer()
        lsa_summary= lsa_summarizer(parser.document,N)
        print(type(lsa_summary))
        path= r"C:\Users\lenovo\Desktop\FYP\backend\Presentation_templates\\"
        path=path+TEMP
        prs = Presentation(path)
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        shapes = slide.shapes
        body_shape = shapes.placeholders[0]
        tf = body_shape.text_frame
        tf.text="SUMMARY"

        for l in lsa_summary:
            print(l)
            print("***************")
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            shapes = slide.shapes
            body_shape = shapes.placeholders[1]
            tf = body_shape.text_frame

            a=str(l)
            if(a.startswith('[')):

                a=a[1:]
            if(a.startswith('\'')):
                    a=a[1:]
            if(a.endswith(']')):

                a=a[:-1]
            if(a.endswith('\'')):
                    a=a[:-1]


            tf.text =a

        prs.save(r'C:\Users\lenovo\Desktop\FYP\backend\Output Presentations\OUTPUT.pptx')
        if(oformat == "pdf") :
            print("reached-1")
            os.system(r'python C:\Users\lenovo\Desktop\FYP\backend\backend\convertion.py')

    #IF THE DOCUMENT UPLOADED IS .DOCX

    elif file_type == "docx":

        print("it is docx")


        text = docx2txt.process(directory, r"C:\Users\lenovo\Desktop\FYP\backend\img")
         #EXTRACT IMAGES FROM DOCX
        Final_theme.generate_ppt(directory,TEMP,N,oformat)



    #IF THE UPLOADED DOCUMENTED IS .PDF

    elif file_type=="pdf":

        dname=r"C:\Users\lenovo\Desktop\FYP\backend\uploads\Convertedpdf.docx"
        #EXTRACTING IMAGES FROM PDF
        doc= fitz.open(directory)
        for i in range(len(doc)):
            for img in doc.getPageImageList(i):
                xref = img[0]
                pix = fitz.Pixmap(doc, xref)
                if pix.n < 5:
                    pix.writePNG(".\img\p%s-%s.png" % (i, xref))
                else:
                    pix1 = fitz.Pixmap(fitz.csRGB, pix)
                    pix1.writePNG(".\img\p%s-%s.png" % (i, xref))
                    pix1 = None
                pix = None
        parse(directory,dname)
        doc.close()
        Final_theme.generate_ppt(dname,TEMP,N,oformat)
