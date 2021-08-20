from docx import Document
from pptx.util import Inches,Pt
from pptx.enum.text import  MSO_AUTO_SIZE
import os
from sumy.summarizers.lsa import LsaSummarizer
from sumy.nlp.tokenizers import Tokenizer
from sumy.parsers.plaintext import PlaintextParser
from pptx import Presentation
import pandas as pd
import sys
#import comtypes.client
import win32com.client

def create_powerpoint(file_name):
	ppt = Presentation()
	save_ppt(ppt, file_name)

def save_ppt(ppt, file_name):
	ppt.save(file_name)

def open_powerpoint(file_name):
	ppt = Presentation(file_name)
	return ppt

def create_slide(ppt, layout):
	return ppt.slides.add_slide(layout)

def add_basic_layouts(ppt):
	ppt_layout = ppt.slide_layouts[0]
	ppt_layout1 = ppt.slide_layouts[1]
	ppt_layout2 = ppt.slide_layouts[2]
	ppt_layout3 = ppt.slide_layouts[3]
	ppt_layout4 = ppt.slide_layouts[4]
	ppt_layout5 = ppt.slide_layouts[5]
	ppt_layout6 = ppt.slide_layouts[6]
	ppt_layout7 = ppt.slide_layouts[7]
	ppt_layout8 = ppt.slide_layouts[8]
	ppt_layout9 = ppt.slide_layouts[9]
	ppt_layout10 = ppt.slide_layouts[10]

	create_slide(ppt, ppt_layout)
	create_slide(ppt, ppt_layout2)
	create_slide(ppt, ppt_layout3)
	create_slide(ppt, ppt_layout4)
	create_slide(ppt, ppt_layout5)
	create_slide(ppt, ppt_layout6)
	create_slide(ppt, ppt_layout7)
	create_slide(ppt, ppt_layout8)
	create_slide(ppt, ppt_layout9)
	create_slide(ppt, ppt_layout10)

def add_blank_layout(ppt):
	ppt_layout = ppt.slide_layouts[6]
	current_slide = create_slide(ppt, ppt_layout)
	return current_slide

def add_text(slide, text, bold, font_size, pos_left, pos_right):
    area_left = Inches(pos_left)
    area_right = Inches(pos_right)
    area_height_width = Inches(1)
    text_area = slide.shapes.add_textbox(area_left, area_right, area_height_width, area_height_width)
    text_frame = text_area.text_frame
    text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    current_paragraph = text_frame.add_paragraph()
    current_paragraph.text = text

    if bold == True:
        current_paragraph.font.bold = True

    current_paragraph.font.size = Pt(font_size)

def add_image(slide, image, pos_left, pos_top, width, *height):
	area_left = Inches(pos_left)
	area_top = Inches(pos_top)
	area_width = Inches(width)

	if len(height) > 0 :
		area_height = Inches(height[0])
		slide.shapes.add_picture(image, area_left, area_top, area_width, area_height)
	else:
		slide.shapes.add_picture(image, area_left, area_top, area_width)

def add_table(slide, nb_rows, nb_columns, pos_left, pos_top, width,dfa,*height):
    df=dfa
    area_left = Inches(pos_left)
    area_top = Inches(pos_top)
    area_width = Inches(width)

    if len(height) > 0 :
        area_height = Inches(height[0])
    else:
        area_height = Inches(width)
    current_table = slide.shapes.add_table(nb_rows, nb_columns, area_left, area_top, area_width,area_height)
    head_list=list(df.columns.values)
    print(head_list)
    len_of_headings=len(head_list)
    for n in head_list:
        print (n)
    k=0
    for m in head_list:
        if(k<len_of_headings):
            current_table.table.cell(0, k).text =m
            k=k+1
    p=1

    for i in range(nb_rows-1):
        for j in range(nb_columns):
            val=df.iat[i,j]
            if(p<=nb_rows):
                current_table.table.cell(p, j).text = val

        p=p+1
    return


def generate_ppt(docx,TEMP,N,oformat):
    print("reached2")
    document = Document(docx)
    headings = []
    texts = []
    para = []
    #SEPERATING HEADINGS AND CONTENT, SENDING CONTENT TO SUMMARIZER
    for paragraph in document.paragraphs:
       for run in paragraph.runs:
            if(run.bold):
                if headings:
                    texts.append(para)
                headings.append(paragraph.text)
                para = []
       if paragraph.style.name == "Normal":
            para.append(paragraph.text)
    if para or len(headings)>len(texts):
        texts.append(texts.append(para))
    ch=len(headings)
    no_of_sentences=N
    for h, t in zip(headings, texts):
        print(h, t)
        print(h)
        parser=PlaintextParser.from_string(t,Tokenizer('english'))
        # creating the summarizer
        lsa_summarizer=LsaSummarizer()
        lsa_summary= lsa_summarizer(parser.document,no_of_sentences)
        print("Summary:")
        print(lsa_summary)
        print("\n\n*****************************")
    ch=len(headings)
    #CREATING A PRESENTATION
    file_name = r"C:\Users\lenovo\Desktop\FYP\backend\Output Presentations\Output.pptx"
    create_powerpoint(file_name)
    ppt = open_powerpoint(file_name)
    directory = r"C:\Users\lenovo\Desktop\FYP\backend\img"
    captions=[]
    path= r"C:\Users\lenovo\Desktop\FYP\backend\Presentation_templates\\"
    path=path+TEMP
    prs = Presentation(path)

    for h, t in zip(headings, texts):

        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)

        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = h
        tf = body_shape.text_frame


        parser=PlaintextParser.from_string(t,Tokenizer('english'))
        lsa_summarizer=LsaSummarizer()
        lsa_summary= lsa_summarizer(parser.document,no_of_sentences)
        b=''
        for l in lsa_summary:
             a=str(l)
             if(a.startswith('[')):

                 a=a[1:]
             if(a.startswith('\'')):
                     a=a[1:]
             if(a.endswith(']')):

                 a=a[:-1]
             if(a.endswith('\'')):
                     a=a[:-1]
             b=b+a


        tf.text =b

    #CREATING THE IMAGE SLIDES
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text(slide1, "IMAGES", True, 60, 5, 4)
    #READING IMAGES FROM SPECIFIED FOLDER AND INSERTING TO PPT
    #REMEMBER TO CLEAR THE IMAGES FOLDER EVERYTIME BEFORE RUNNING THE PROGRAM
    for filename in os.listdir(directory):
        if filename.endswith(".jpg") or filename.endswith(".png") or filename.endswith(".jpeg"):
            print(os.path.join(directory, filename))
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_image(slide, os.path.join(directory, filename),3,2,4,4)
        else:
            continue
    for f in os.listdir(directory):
        os.remove(os.path.join(directory, f))


    #ADDING TABLE TO PPT

    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text(slide2, "TABLES", True, 60, 5, 4)
    table_count = len(document.tables)

    for i in range(table_count):
        table = document.tables[i]
        data = []
        keys = None
        for i, row in enumerate(table.rows):
            text = (cell.text for cell in row.cells)
            if i == 0:
                keys = tuple(text)
                continue
            row_data = dict(zip(keys, text))
            data.append(row_data)
        df = pd.DataFrame(data)
        print(df)
        nbr=df.shape[0]
        nbc=df.shape[1]
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_table(slide, nbr+1, nbc, 1, 1, 8, df,2)

    #SAVING THE PPT
    save_ppt(ppt, file_name)
    prs.save(r'C:\Users\lenovo\Desktop\FYP\backend\Output Presentations\OUTPUT.pptx')

    if(oformat == "pdf") :
        #inputFileName=r'C:\Users\lenovo\Desktop\FYP\backend\Output Presentations\OUTPUT.pptx'
        #outputFileName= r'C:\Users\lenovo\Desktop\FYP\backend\Output Presentations\Output.pdf'
        #convertion.convert(inputFileName)
        print("reached-1")
        os.system(r'python C:\Users\lenovo\Desktop\FYP\backend\backend\convertion.py')
